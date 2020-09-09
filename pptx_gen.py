# %%
from pptx import Presentation

pptx_path = 'templates/template.pptx'

prs = Presentation(pptx_path)

# %%

# prs.slides[3].shapes[4].table.cell(0,0).text = 'this is my test'


# # %%
# slide = prs.slides[3]

# # %%

# for shape in slide.shapes:
#     if shape.shape_type == 'table':
#         print(shape)
# %%
# shapes = []

# for shape in slide.shapes:
#     if shape.has_table:
#         shapes.append(shape)

# shapes
# shape_to_change = shapes[1]
# shape_to_change.table.cell(0,0).text = "did this work" # this does work
# %%

# %%
def replace_text(replacements, shapes):
    """Takes dict of {match: replacement, ... } and replaces all matches.
    Currently not implemented for charts or graphics.
    """
    for shape in shapes:
        for match, replacement in replacements.items():
            if shape.has_text_frame:
                if (shape.text.find(match)) != -1:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        whole_text = "".join(run.text for run in paragraph.runs)
                        whole_text = whole_text.replace(str(match), str(replacement))
                        for idx, run in enumerate(paragraph.runs):
                            if idx != 0:
                                p = paragraph._p
                                p.remove(run._r)
                        if(not(not paragraph.runs)):
                            paragraph.runs[0].text = whole_text


# %%
# To get shapes in your slides
slides = [slide for slide in prs.slides]
shapes = []
for slide in slides:
    for shape in slide.shapes:
        shapes.append(shape)
replaces = {'Key Themes:': 'Key Farts:',
            'Creating Digital Fluency': 'farts'}
replace_text(replaces, shapes)
prs.save('test.pptx')


# %%
prs.save('test.pptx')

# %%
prs.slides[3].shapes[0].table.cell(0,0).text = 'this is shape 0, (0,0)'
prs.slides[3].shapes[0].table.cell(1,0).text = 'this is shape 0, (1,0)'
prs.slides[3].shapes[0].table.cell(0,1).text = 'this is shape 0, (0,1)'
prs.slides[3].shapes[0].table.cell(1,1).text = 'this is shape 0, (1,1)'

prs.slides[3].shapes[1].table.cell(0,0).text = 'this is shape 1, (0,0)'
prs.slides[3].shapes[1].table.cell(1,0).text = 'this is shape 1, (1,0)'
prs.slides[3].shapes[1].table.cell(0,1).text = 'this is shape 1, (0,1)'
prs.slides[3].shapes[1].table.cell(1,1).text = 'this is shape 1, (1,1)'

prs.save('test.pptx')

# %%
len(prs.slides[3].shapes)
# %%
p = prs.slides[3].shapes[0].table.cell(1,2).text_frame.paragraphs[0]
# %%
p.text = 'this is my test'

prs.save('test.pptx')

# %%
p.font.size
# %%
p_2 = prs.slides[3].shapes[0].table.cell(1,3).text_frame.paragraphs[0]
p_2.font.size

# %%
shape = 


